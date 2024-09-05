import streamlit as st
from sentence_transformers import SentenceTransformer
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import fitz  # PyMuPDF for PDF processing
import docx
import pptx
import openpyxl

st.title("Conversational Chatbot")
st.sidebar.title("Upload Documents")

# Upload documents
uploaded_files = st.sidebar.file_uploader(
    "Upload your documents", 
    type=["txt", "pdf", "docx", "pptx", "xlsx"], 
    accept_multiple_files=True
)

# Load the Sentence Transformer model
model = SentenceTransformer('all-MiniLM-L6-v2')
st.write("Embedding model loaded successfully!")

# Initialize an empty list to store document texts
documents = []

def extract_text_from_file(uploaded_file):
    text = ""
    try:
        if uploaded_file.type == "text/plain":
            text = uploaded_file.read().decode("utf-8")
        elif uploaded_file.type == "application/pdf":
            pdf_document = fitz.open(stream=uploaded_file.read())
            text = "".join([page.get_text() for page in pdf_document])
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = pptx.Presentation(uploaded_file)
            text = "\n".join(
                [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
            )
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            excel_file = openpyxl.load_workbook(uploaded_file)
            text = "\n".join(
                [" ".join([str(cell) for cell in row]) for sheet in excel_file.sheetnames for row in excel_file[sheet].iter_rows(values_only=True)]
            )
    except Exception as e:
        st.write(f"Error processing file {uploaded_file.name}: {e}")
    return text

# Process uploaded files
if uploaded_files:
    for uploaded_file in uploaded_files:
        documents.append(extract_text_from_file(uploaded_file))

    st.write(f"{len(documents)} document(s) uploaded successfully!")

# Generate embeddings for each document if documents are available
embeddings = None
if documents:
    embeddings = model.encode(documents)
    st.write("Document embeddings generated successfully!")

# User input for queries
user_query = st.text_input("Enter your question:")

# Handle user queries
if user_query:
    try:
        if embeddings is not None:
            query_embedding = model.encode([user_query])
            similarities = cosine_similarity(query_embedding, embeddings)
            most_similar_doc_index = np.argmax(similarities)
            st.write("Chatbot response:")
            st.write(documents[most_similar_doc_index])
        else:
            # Use predefined responses if no documents are available
            predefined_responses = {
                "hello": "Hi there! How can I assist you today?",
                "problem": "Tell me more about the problem you're facing.",
                "thank you": "You're welcome! If you have any more questions, feel free to ask.",
                "help": "I'm here to help! What do you need assistance with?",
            }
            response = predefined_responses.get(user_query.lower(), "Sorry, I didn't understand that. Can you please rephrase?")
            st.write("Chatbot response:")
            st.write(response)
    except Exception as e:
        st.write(f"An error occurred: {e}")
